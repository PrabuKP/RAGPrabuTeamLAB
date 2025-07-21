import os
import uuid
from typing import List, Dict, Any

from fastapi import FastAPI, UploadFile, File, HTTPException, Query
from fastapi.staticfiles import StaticFiles
from fastapi.responses    import FileResponse
from PyPDF2        import PdfReader
from docx          import Document as DocxDocument
from pptx          import Presentation
from elasticsearch  import Elasticsearch, ApiError

from llama_index.text_splitter import TokenTextSplitter

app = FastAPI(title="RAGPrabu with UI & File Upload", version="1.0.0")

# mount folder static untuk UI
app.mount("/static", StaticFiles(directory="static"), name="static")

@app.get("/")
def root():
    return FileResponse("static/index.html")

# Elasticsearch config
ES_HOST = os.getenv("ES_HOST", "http://localhost:9200")
ES_INDEX = "rag_docs"
es: Elasticsearch

# folder Data/ untuk UI
DATA_DIR = "Data"

# llama-index splitter
splitter = TokenTextSplitter(chunk_size=300, chunk_overlap=50)

@app.on_event("startup")
def startup_event():
    global es
    es = Elasticsearch(hosts=[ES_HOST])
    try:
        if not es.indices.exists(index=ES_INDEX):
            es.indices.create(
                index=ES_INDEX,
                body={
                    "settings": {"similarity": {"default": {"type": "BM25"}}},
                    "mappings": {
                        "properties": {
                            "doc_id": {"type": "keyword"},
                            "chunk":  {"type": "text"}
                        }
                    }
                }
            )
    except ApiError as e:
        print(f"⚠️  Warning creating index: {e}")

@app.get("/health")
def health() -> Dict[str, str]:
    return {"status": "ok"}

def extract_text(path: str, ext: str) -> str:
    if ext == "pdf":
        reader = PdfReader(path)
        return "\n".join(p.extract_text() or "" for p in reader.pages)
    if ext in ("docx", "doc"):
        doc = DocxDocument(path)
        return "\n".join(p.text for p in doc.paragraphs)
    if ext in ("pptx", "ppt"):
        prs = Presentation(path)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    t = shape.text.strip()
                    if t: texts.append(t)
        return "\n".join(texts)
    # txt/json/py
    with open(path, "r", encoding="utf-8") as f:
        return f.read()

@app.get("/data-files")
def data_files() -> Dict[str, List[str]]:
    """List nama file di folder Data/ untuk UI."""
    files = []
    for fn in os.listdir(DATA_DIR):
        full = os.path.join(DATA_DIR, fn)
        if os.path.isfile(full):
            files.append(fn)
    return {"files": files}


@app.post("/upload-from-data")
def upload_from_data(filename: str = Query(..., description="Nama file di Data/")):
    """
    Upload & index file yang sudah ada di Data/.
    Dipanggil oleh UI (index.html).
    """
    path = os.path.join(DATA_DIR, filename)
    if not os.path.isfile(path):
        raise HTTPException(404, f"File not found: {filename}")

    # sama seperti /upload: extract, chunk, index
    doc_id = str(uuid.uuid4())
    ext = filename.rsplit(".", 1)[1].lower()
    raw = extract_text(path, ext)

    chunks = splitter.split_text(raw)
    for i, chunk in enumerate(chunks):
        try:
            es.index(
                index=ES_INDEX,
                id=f"{doc_id}_{i}",
                body={"doc_id": doc_id, "chunk": chunk}
            )
        except ApiError as e:
            print(f"⚠️  failed to index chunk {i}: {e}")

    es.indices.refresh(index=ES_INDEX)
    return {"document_id": doc_id, "chunk_count": len(chunks)}


@app.post("/upload")
async def upload_file(file: UploadFile = File(...)):
    """
    Upload file arbitrary via multipart-form.
    Dipanggil oleh client.py atau shell script.
    """
    name, ext = file.filename.rsplit(".", 1)
    ext = ext.lower()
    if ext not in ("pptx","ppt","docx","doc","pdf","txt","json","py"):
        raise HTTPException(415, f"Unsupported file type: .{ext}")

    # simpan sementara
    doc_id = str(uuid.uuid4())
    tmp = f"/tmp/{doc_id}.{ext}"
    with open(tmp, "wb") as f:
        f.write(await file.read())

    # extract & index
    raw = extract_text(tmp, ext)
    os.remove(tmp)

    chunks = splitter.split_text(raw)
    for i, chunk in enumerate(chunks):
        try:
            es.index(
                index=ES_INDEX,
                id=f"{doc_id}_{i}",
                body={"doc_id": doc_id, "chunk": chunk}
            )
        except ApiError as e:
            print(f"⚠️  failed to index chunk {i}: {e}")

    es.indices.refresh(index=ES_INDEX)
    return {"document_id": doc_id, "chunk_count": len(chunks)}


@app.get("/retrieve")
def retrieve(
    document_id: str,
    question:    str,
    top_k:       int = 5
) -> Dict[str, Any]:
    """
    Retrieve top-k BM25 chunks untuk document_id + question.
    """
    body = {
        "size": top_k,
        "query": {
            "bool": {
                "must": [
                    {"term":  {"doc_id": document_id}},
                    {"match": {"chunk": {"query": question}}}
                ]
            }
        }
    }
    try:
        res = es.search(index=ES_INDEX, body=body)
    except ApiError as e:
        raise HTTPException(500, f"Search error: {e}")

    frags = [
        {"chunk": h["_source"]["chunk"], "score": h["_score"]}
        for h in res["hits"]["hits"]
    ]
    return {"fragments": frags}
