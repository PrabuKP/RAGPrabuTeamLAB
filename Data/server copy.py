import os
import uuid
from typing import List, Dict, Any
from fastapi import FastAPI, UploadFile, File, HTTPException
from PyPDF2 import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation
from elasticsearch import Elasticsearch, ApiError
from llama_index.text_splitter import TokenTextSplitter

app = FastAPI(title="RAGPrabu API w/ llama-index Chunking", version="0.1.0")

# Elasticsearch config
ES_HOST = os.getenv("ES_HOST", "http://localhost:9200")
ES_INDEX = "rag_docs"
es: Elasticsearch

# Inisialisasi splitter: 200 token per chunk, overlap 50 token
splitter = TokenTextSplitter(chunk_size=200, chunk_overlap=50)

@app.on_event("startup")
def on_startup():
    global es
    es = Elasticsearch(hosts=[ES_HOST])
    try:
        if not es.indices.exists(index=ES_INDEX):
            es.indices.create(
                index=ES_INDEX,
                body={
                    "settings": {
                        "similarity": {
                            "default": {"type": "BM25"}
                        }
                    },
                    "mappings": {
                        "properties": {
                            "doc_id": {"type": "keyword"},
                            "chunk":  {"type": "text"}
                        }
                    }
                }
            )
    except ApiError as e:
        print(f"⚠️  Warning creating index: {e}")

@app.get("/health")
def health() -> Dict[str, str]:
    return {"status": "ok"}

def extract_text(path: str, ext: str) -> str:
    if ext == "pdf":
        reader = PdfReader(path)
        return "\n".join(p.extract_text() or "" for p in reader.pages)
    if ext in ("docx","doc"):
        doc = DocxDocument(path)
        return "\n".join(p.text for p in doc.paragraphs)
    if ext in ("pptx","ppt"):
        prs = Presentation(path)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    t = shape.text.strip()
                    if t:
                        texts.append(t)
        return "\n".join(texts)
    # txt, json, py
    with open(path, "r", encoding="utf-8") as f:
        return f.read()

@app.post("/upload")
async def upload(file: UploadFile = File(...)) -> Dict[str, Any]:
    # Validasi ekstensi
    name, ext = file.filename.rsplit(".", 1)
    ext = ext.lower()
    if ext not in ("pptx","ppt","docx","doc","pdf","txt","json","py"):
        raise HTTPException(415, f"Unsupported file type: .{ext}")

    # Simpan sementara & extract text
    doc_id   = str(uuid.uuid4())
    tmp_path = f"/tmp/{doc_id}.{ext}"
    with open(tmp_path, "wb") as f:
        f.write(await file.read())
    raw = extract_text(tmp_path, ext)
    os.remove(tmp_path)

    # Chunking lewat llama-index TokenTextSplitter
    frags = splitter.split_text(raw)

    # Index tiap chunk ke Elasticsearch
    for i, chunk in enumerate(frags):
        try:
            es.index(
                index=ES_INDEX,
                id=f"{doc_id}_{i}",
                body={"doc_id": doc_id, "chunk": chunk}
            )
        except ApiError as e:
            print(f"⚠️  Failed to index chunk {i}: {e}")

    # Refresh agar segera searchable
    try:
        es.indices.refresh(index=ES_INDEX)
    except ApiError:
        pass

    return {"document_id": doc_id, "chunk_count": len(frags)}

@app.get("/retrieve")
def retrieve(
    document_id: str,
    question:    str,
    top_k:       int = 5
) -> Dict[str, List[Dict[str, Any]]]:
    # BM25 retrieval: must term doc_id + must match chunk
    body = {
        "size": top_k,
        "query": {
            "bool": {
                "must": [
                    {"term":  {"doc_id": document_id}},
                    {"match": {"chunk": {"query": question}}}
                ]
            }
        }
    }
    try:
        res = es.search(index=ES_INDEX, body=body)
    except ApiError as e:
        raise HTTPException(500, f"Search error: {e}")

    return {
        "fragments": [
            {"chunk": hit["_source"]["chunk"], "score": hit["_score"]}
            for hit in res["hits"]["hits"]
        ]
    }
