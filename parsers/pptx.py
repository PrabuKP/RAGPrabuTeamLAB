from pptx import Presentation

def extract_text(path: str) -> str:
    prs = Presentation(path)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                txt = shape.text.strip()
                if txt:
                    text_runs.append(txt)
    return "\n".join(text_runs)
